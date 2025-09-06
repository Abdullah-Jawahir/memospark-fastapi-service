import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation
import io
from fastapi import HTTPException
from .utils import clean_text
from .logger import logger

def extract_text_from_pdf(file_content: bytes) -> str:
    """Extract text from PDF file."""
    try:
        doc = fitz.open(stream=file_content, filetype="pdf")
        text = ""
        for page in doc:
            page_text = page.get_text()
            text += page_text + "\n"
        doc.close()
        
        # Clean the extracted text
        text = clean_text(text)
        
        logger.info(f"Successfully extracted {len(text)} characters from PDF")
        return text
    except Exception as e:
        error_msg = f"Error extracting text from PDF: {str(e)}"
        logger.error(error_msg, exc_info=True)
        raise HTTPException(status_code=400, detail="Error processing PDF file")

def extract_text_from_docx(file_content: bytes) -> str:
    """Extract text from DOCX file."""
    try:
        doc = DocxDocument(io.BytesIO(file_content))
        text = ""
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text += paragraph.text.strip() + "\n"
        
        # Clean the extracted text
        text = clean_text(text)
        
        logger.info(f"Successfully extracted {len(text)} characters from DOCX")
        return text
    except Exception as e:
        error_msg = f"Error extracting text from DOCX: {str(e)}"
        logger.error(error_msg, exc_info=True)
        raise HTTPException(status_code=400, detail="Error processing DOCX file")

def extract_text_from_pptx(file_content: bytes) -> str:
    """Extract text from PPTX file."""
    try:
        prs = Presentation(io.BytesIO(file_content))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text.strip() + "\n"
        
        # Clean the extracted text
        text = clean_text(text)
        
        logger.info(f"Successfully extracted {len(text)} characters from PPTX")
        return text
    except Exception as e:
        error_msg = f"Error extracting text from PPTX: {str(e)}"
        logger.error(error_msg, exc_info=True)
        raise HTTPException(status_code=400, detail="Error processing PPTX file")

def extract_text_from_txt(file_content: bytes) -> str:
    """Extract text from TXT file."""
    try:
        # Try to decode the text with common encodings
        for encoding in ['utf-8', 'utf-8-sig', 'latin-1', 'cp1252']:
            try:
                text = file_content.decode(encoding)
                break
            except UnicodeDecodeError:
                continue
        else:
            # If all encodings fail, use utf-8 with error handling
            text = file_content.decode('utf-8', errors='replace')
        
        # Clean the extracted text
        text = clean_text(text)
        
        logger.info(f"Successfully extracted {len(text)} characters from TXT")
        return text
    except Exception as e:
        error_msg = f"Error extracting text from TXT: {str(e)}"
        logger.error(error_msg, exc_info=True)
        raise HTTPException(status_code=400, detail="Error processing TXT file")

def extract_text_from_file(file_content: bytes, file_extension: str) -> str:
    """Extract text from file based on its extension."""
    file_extension = file_extension.lower()
    
    if file_extension == "pdf":
        return extract_text_from_pdf(file_content)
    elif file_extension == "docx":
        return extract_text_from_docx(file_content)
    elif file_extension == "pptx":
        return extract_text_from_pptx(file_content)
    elif file_extension == "txt":
        return extract_text_from_txt(file_content)
    else:
        raise HTTPException(status_code=400, detail="Unsupported file type") 