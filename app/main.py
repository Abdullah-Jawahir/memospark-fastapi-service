from fastapi import FastAPI, UploadFile, File, HTTPException, Form, Depends
from fastapi.middleware.cors import CORSMiddleware
import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation
import io
from transformers import pipeline, AutoTokenizer, AutoModelForSeq2SeqLM, AutoModelForCausalLM
import torch
from typing import List, Dict, Any, Optional
import logging
import os
from pathlib import Path
from datetime import datetime
import re
import random
from fastapi import Request
from fastapi import Query
from fastapi import APIRouter
import json

# Configure logging with file handler
PROJECT_ROOT = Path(__file__).parent.parent
LOG_DIR = PROJECT_ROOT / "logs"
LOG_DIR.mkdir(exist_ok=True)

# Create a logger
logger = logging.getLogger(__name__)
logger.setLevel(logging.DEBUG)

# Create file handler
log_file = LOG_DIR / f"fastapi_errors_{datetime.now().strftime('%Y%m%d')}.log"
file_handler = logging.FileHandler(log_file)
file_handler.setLevel(logging.ERROR)

# Create console handler
console_handler = logging.StreamHandler()
console_handler.setLevel(logging.INFO)

# Create formatter
formatter = logging.Formatter('%(asctime)s - %(name)s - %(levelname)s - %(message)s')
file_handler.setFormatter(formatter)
console_handler.setFormatter(formatter)

# Add handlers to logger
logger.addHandler(file_handler)
logger.addHandler(console_handler)

# Set custom cache directory for Hugging Face models
CACHE_DIR = PROJECT_ROOT / "model_cache"
CACHE_DIR.mkdir(exist_ok=True)

os.environ['TRANSFORMERS_CACHE'] = str(CACHE_DIR)
os.environ['HF_HOME'] = str(CACHE_DIR)
os.environ['HF_DATASETS_CACHE'] = str(CACHE_DIR)

app = FastAPI()

# Configure CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # In production, replace with specific origins
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Initialize multiple models for different tasks
MODEL_CONFIGS = {
    "question_generation": {
        "name": "microsoft/DialoGPT-medium",  # Better for conversational Q&A
        "alternative": "google/flan-t5-base"  # Larger T5 model
    },
    "text_generation": {
        "name": "microsoft/DialoGPT-small",
        "alternative": "google/flan-t5-large"
    }
}

# Try to load the best available model
def load_best_model():
    """Load the best available model for text generation tasks."""
    models_to_try = [
        "google/flan-t5-base",  # Better than small version
        "microsoft/DialoGPT-medium",
        "google/flan-t5-small",  # Fallback
    ]
    
    for model_name in models_to_try:
        try:
            logger.info(f"Attempting to load model: {model_name}")
            tokenizer = AutoTokenizer.from_pretrained(model_name, cache_dir=str(CACHE_DIR))
            
            if "flan-t5" in model_name:
                model = AutoModelForSeq2SeqLM.from_pretrained(model_name, cache_dir=str(CACHE_DIR))
            else:
                model = AutoModelForCausalLM.from_pretrained(model_name, cache_dir=str(CACHE_DIR))
            
            logger.info(f"Successfully loaded model: {model_name}")
            return tokenizer, model, model_name
        except Exception as e:
            logger.warning(f"Failed to load {model_name}: {str(e)}")
            continue
    
    raise Exception("Failed to load any suitable model")

try:
    tokenizer, model, current_model_name = load_best_model()
    logger.info(f"Using model: {current_model_name}")
except Exception as e:
    logger.error(f"Failed to load any model: {str(e)}")
    raise

# Initialize question generation pipeline
try:
    question_generator = pipeline(
        "text2text-generation" if "flan-t5" in current_model_name else "text-generation",
        model=model,
        tokenizer=tokenizer,
        device=-1  # Use CPU
    )
    logger.info("Question generation pipeline loaded successfully")
except Exception as e:
    logger.error(f"Failed to load question generation pipeline: {str(e)}")
    question_generator = None

def clean_text(text: str) -> str:
    """Clean and normalize extracted text."""
    # Remove extra whitespace and normalize
    text = re.sub(r'\s+', ' ', text)
    text = text.strip()
    
    # Remove common PDF artifacts
    text = re.sub(r'[^\w\s\.\,\;\:\!\?\-\(\)\[\]\'\"]+', ' ', text)
    
    # Fix common OCR issues
    text = re.sub(r'(\w)\|(\w)', r'\1l\2', text)  # Fix common OCR 'l' -> '|' issue
    text = re.sub(r'(\w)0(\w)', r'\1o\2', text)   # Fix common OCR 'o' -> '0' issue
    
    return text

def extract_key_concepts(text: str, max_concepts: int = 10) -> List[str]:
    """Extract key concepts from text for better question generation."""
    # Simple keyword extraction based on capitalization and frequency
    words = re.findall(r'\b[A-Z][a-z]+(?:\s+[A-Z][a-z]+)*\b', text)
    
    # Count word frequency
    word_freq = {}
    for word in words:
        word_freq[word] = word_freq.get(word, 0) + 1
    
    # Sort by frequency and return top concepts
    key_concepts = sorted(word_freq.items(), key=lambda x: x[1], reverse=True)
    return [concept[0] for concept in key_concepts[:max_concepts]]

def extract_text_from_pdf(file_content: bytes) -> str:
    """Extract text from PDF file."""
    try:
        doc = fitz.open(stream=file_content, filetype="pdf")
        text = ""
        for page in doc:
            page_text = page.get_text()
            text += page_text + "\n"
        doc.close()
        
        # Clean the extracted text
        text = clean_text(text)
        
        logger.info(f"Successfully extracted {len(text)} characters from PDF")
        return text
    except Exception as e:
        error_msg = f"Error extracting text from PDF: {str(e)}"
        logger.error(error_msg, exc_info=True)
        raise HTTPException(status_code=400, detail="Error processing PDF file")

def extract_text_from_docx(file_content: bytes) -> str:
    """Extract text from DOCX file."""
    try:
        doc = DocxDocument(io.BytesIO(file_content))
        text = ""
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text += paragraph.text.strip() + "\n"
        
        # Clean the extracted text
        text = clean_text(text)
        
        logger.info(f"Successfully extracted {len(text)} characters from DOCX")
        return text
    except Exception as e:
        error_msg = f"Error extracting text from DOCX: {str(e)}"
        logger.error(error_msg, exc_info=True)
        raise HTTPException(status_code=400, detail="Error processing DOCX file")

def extract_text_from_pptx(file_content: bytes) -> str:
    """Extract text from PPTX file."""
    try:
        prs = Presentation(io.BytesIO(file_content))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text.strip() + "\n"
        
        # Clean the extracted text
        text = clean_text(text)
        
        logger.info(f"Successfully extracted {len(text)} characters from PPTX")
        return text
    except Exception as e:
        error_msg = f"Error extracting text from PPTX: {str(e)}"
        logger.error(error_msg, exc_info=True)
        raise HTTPException(status_code=400, detail="Error processing PPTX file")

def generate_with_model(prompt: str, max_length: int = 200) -> str:
    """Generate text using the loaded model."""
    try:
        inputs = tokenizer(
            prompt,
            return_tensors="pt",
            max_length=1024,
            truncation=True,
            padding=True
        )
        
        with torch.no_grad():
            if "flan-t5" in current_model_name:
                outputs = model.generate(
                    inputs["input_ids"],
                    max_length=max_length,
                    num_return_sequences=1,
                    temperature=0.7,
                    top_p=0.9,
                    do_sample=True,
                    pad_token_id=tokenizer.eos_token_id,
                    eos_token_id=tokenizer.eos_token_id,
                    no_repeat_ngram_size=2
                )
            else:
                outputs = model.generate(
                    inputs["input_ids"],
                    max_length=len(inputs["input_ids"][0]) + max_length,
                    num_return_sequences=1,
                    temperature=0.7,
                    top_p=0.9,
                    do_sample=True,
                    pad_token_id=tokenizer.eos_token_id,
                    eos_token_id=tokenizer.eos_token_id,
                    no_repeat_ngram_size=2
                )
        
        generated_text = tokenizer.decode(outputs[0], skip_special_tokens=True)
        
        # For non-T5 models, remove the input prompt from output
        if "flan-t5" not in current_model_name:
            generated_text = generated_text[len(prompt):].strip()
        
        return generated_text
    except Exception as e:
        logger.error(f"Error in model generation: {str(e)}")
        return ""

def generate_flashcards(text: str, language: str = "en", difficulty: str = "beginner") -> List[Dict[str, Any]]:
    """Generate flashcards from text using improved prompting and parsing."""
    try:
        text = text.strip()
        if len(text) < 50:
            logger.warning("Text is too short to generate meaningful flashcards")
            return []
        
        # Extract key concepts for better question generation
        key_concepts = extract_key_concepts(text)
        
        # Split text into meaningful chunks
        sentences = [s.strip() for s in text.split('.') if len(s.strip()) > 30]
        chunks = []
        current_chunk = ""
        
        for sentence in sentences:
            if len(current_chunk + sentence) > 600:
                if current_chunk:
                    chunks.append(current_chunk.strip())
                current_chunk = sentence
            else:
                current_chunk += ". " + sentence if current_chunk else sentence
        
        if current_chunk:
            chunks.append(current_chunk.strip())
        
        flashcards = []
        
        # Language-specific prompts
        prompts = {
            "en": "Create a clear question and answer from this text:\n\nText: {}\n\nQuestion:",
            "si": "මෙම පෙළෙන් පැහැදිලි ප්‍රශ්නයක් සහ පිළිතුරක් සාදන්න:\n\nපෙළ: {}\n\nප්‍රශ්නය:",
            "ta": "இந்த உரையிலிருந்து தெளிவான கேள்வி மற்றும் பதிலை உருவாக்கவும்:\n\nஉரை: {}\n\nகேள்வி:"
        }
        
        base_prompt = prompts.get(language, prompts["en"])
        
        for i, chunk in enumerate(chunks[:5]):  # Limit to 5 chunks
            try:
                # Use key concepts to generate focused questions
                if key_concepts and i < len(key_concepts):
                    concept = key_concepts[i]
                    focused_prompt = f"{base_prompt.format(chunk)} What is {concept}?"
                else:
                    focused_prompt = base_prompt.format(chunk)
                
                generated = generate_with_model(focused_prompt, max_length=150)
                
                if generated and len(generated.strip()) > 10:
                    # Try to parse question and answer
                    lines = generated.strip().split('\n')
                    question = ""
                    answer = ""
                    
                    for line in lines:
                        line = line.strip()
                        if line.startswith(('Question:', 'ප්‍රශ්නය:', 'கேள்வி:')):
                            question = line.split(':', 1)[1].strip()
                        elif line.startswith(('Answer:', 'පිළිතුර:', 'பதில்:')):
                            answer = line.split(':', 1)[1].strip()
                        elif not question and len(line) > 10 and line.endswith('?'):
                            question = line
                        elif question and not answer and len(line) > 10:
                            answer = line
                    
                    # If parsing failed, create from the chunk
                    if not question:
                        if language == "si":
                            question = f"මෙම කරුණ ගැන කුමක් කිව හැකිද: {chunk[:100]}...?"
                        elif language == "ta":
                            question = f"இந்த விஷயத்தைப் பற்றி என்ன சொல்லலாம்: {chunk[:100]}...?"
                        else:
                            question = f"What can you tell me about: {chunk[:100]}...?"
                    
                    if not answer:
                        answer = chunk[:200] + "..." if len(chunk) > 200 else chunk
                    
                    # Clean up and validate
                    question = re.sub(r'\s+', ' ', question).strip()
                    answer = re.sub(r'\s+', ' ', answer).strip()
                    
                    if len(question) > 10 and len(answer) > 10:
                        flashcards.append({
                            "question": question,
                            "answer": answer,
                            "type": "Q&A",
                            "difficulty": difficulty
                        })
                        
            except Exception as e:
                logger.warning(f"Error generating flashcard for chunk {i}: {str(e)}")
                continue
        
        # If model failed, generate rule-based flashcards
        if not flashcards:
            flashcards = generate_rule_based_flashcards(text, language, difficulty)
        
        logger.info(f"Generated {len(flashcards)} flashcards")
        return flashcards[:10]  # Limit to 10
        
    except Exception as e:
        logger.error(f"Error generating flashcards: {str(e)}")
        return generate_rule_based_flashcards(text, language, difficulty)

def generate_rule_based_flashcards(text: str, language: str, difficulty: str) -> List[Dict[str, Any]]:
    """Generate flashcards using rule-based approach when model fails."""
    flashcards = []
    sentences = [s.strip() for s in text.split('.') if len(s.strip()) > 20]
    key_concepts = extract_key_concepts(text)
    
    for i, sentence in enumerate(sentences[:8]):
        try:
            if language == "si":
                question = f"'{sentence[:50]}...' යන්නෙන් අදහස් කරන්නේ කුමක්ද?"
            elif language == "ta":
                question = f"'{sentence[:50]}...' என்பதன் பொருள் என்ன?"
            else:
                question = f"What does '{sentence[:50]}...' refer to?"
            
            flashcards.append({
                "question": question,
                "answer": sentence,
                "type": "Q&A",
                "difficulty": difficulty
            })
        except Exception as e:
            continue
    
    return flashcards

def generate_quizzes(text: str, language: str = "en", difficulty: str = "beginner") -> List[Dict[str, Any]]:
    """Generate multiple choice quizzes with improved accuracy."""
    try:
        text = text.strip()
        if len(text) < 100:
            return []
        
        sentences = [s.strip() for s in text.split('.') if len(s.strip()) > 30]
        key_concepts = extract_key_concepts(text)
        quizzes = []
        
        # Generate quiz prompts
        quiz_prompts = {
            "en": "Create a multiple choice question with 4 options from this text:\n\nText: {}\n\nQuestion: ",
            "si": "මෙම පෙළෙන් විකල්ප 4ක් සහිත බහුවරණ ප්‍රශ්නයක් සාදන්න:\n\nපෙළ: {}\n\nප්‍රශ්නය: ",
            "ta": "இந்த உரையிலிருந்து 4 விருப்பங்களுடன் பல தேர்வு கேள்வியை உருவாக்கவும்:\n\nஉரை: {}\n\nகேள்வி: "
        }
        
        base_prompt = quiz_prompts.get(language, quiz_prompts["en"])
        
        for i, sentence in enumerate(sentences[:6]):
            try:
                if len(sentence.split()) < 8:
                    continue
                
                # Generate question using model
                prompt = base_prompt.format(sentence)
                generated = generate_with_model(prompt, max_length=200)
                
                # Create options: 1 correct, 3 distractors
                correct_answer = sentence
                distractors = []
                
                # Get other sentences as distractors
                other_sentences = [s for j, s in enumerate(sentences) if j != i and len(s.split()) > 5]
                random.shuffle(other_sentences)
                distractors = other_sentences[:3]
                
                # If not enough distractors, create some
                while len(distractors) < 3:
                    if language == "si":
                        distractors.append(f"වෙනත් සාධාරණ තොරතුර {len(distractors) + 1}")
                    elif language == "ta":
                        distractors.append(f"பிற பொதுவான தகவல் {len(distractors) + 1}")
                    else:
                        distractors.append(f"Other general information {len(distractors) + 1}")
                
                options = distractors + [correct_answer]
                random.shuffle(options)
                
                # Generate question
                words = sentence.split()
                if language == "si":
                    question = f"පහත සඳහන් කුමන කාරණය සත්‍ය ද? {' '.join(words[:8])}... ගැන"
                elif language == "ta":
                    question = f"பின்வருவனவற்றில் எது உண்மை? {' '.join(words[:8])}... பற்றி"
                else:
                    question = f"Which of the following is true about: {' '.join(words[:8])}...?"
                
                quizzes.append({
                    "question": question,
                    "options": options,
                    "answer": correct_answer,
                    "correct_answer_option": correct_answer,
                    "type": "quiz",
                    "difficulty": difficulty
                })
                
            except Exception as e:
                logger.warning(f"Error generating quiz {i}: {str(e)}")
                continue
        
        logger.info(f"Generated {len(quizzes)} quizzes")
        return quizzes[:8]  # Limit to 8
        
    except Exception as e:
        logger.error(f"Error generating quizzes: {str(e)}")
        return []

def generate_exercises(text: str, language: str = "en", difficulty: str = "beginner") -> List[Dict[str, Any]]:
    """Generate various types of exercises from text."""
    try:
        text = text.strip()
        if len(text) < 50:
            return []
        
        exercises = []
        sentences = [s.strip() for s in text.split('.') if len(s.strip()) > 20]
        key_concepts = extract_key_concepts(text)
        
        # Exercise Type 1: Fill in the blanks
        for i, sentence in enumerate(sentences[:3]):
            try:
                words = sentence.split()
                if len(words) > 8:
                    # Remove important words (nouns, adjectives)
                    important_words = [w for w in words if len(w) > 4 and w.isalpha()]
                    if important_words:
                        word_to_blank = random.choice(important_words)
                        exercise_text = sentence.replace(word_to_blank, "______")
                        
                        if language == "si":
                            instruction = "පහත වාක්‍යයේ හිස් තැන පුරවන්න:"
                        elif language == "ta":
                            instruction = "பின்வரும் வாக््यத்தில் காली இடத்தை நிரப்பவும்:"
                        else:
                            instruction = "Fill in the blank in the following sentence:"
                        
                        exercises.append({
                            "type": "fill_blank",
                            "instruction": instruction,
                            "exercise_text": exercise_text,
                            "answer": word_to_blank,
                            "difficulty": difficulty
                        })
            except Exception as e:
                continue
        
        # Exercise Type 2: True/False
        for i, sentence in enumerate(sentences[3:6]):
            try:
                if language == "si":
                    instruction = "පහත ප්‍රකාශනය සත්‍ය ද අසත්‍ය ද?"
                elif language == "ta":
                    instruction = "பின்வரும் கூற்று உண்மையா பொய்யா?"
                else:
                    instruction = "Is the following statement true or false?"
                
                exercises.append({
                    "type": "true_false",
                    "instruction": instruction,
                    "exercise_text": sentence,
                    "answer": "True",  # Assuming text content is factual
                    "difficulty": difficulty
                })
            except Exception as e:
                continue
        
        # Exercise Type 3: Short Answer
        for concept in key_concepts[:2]:
            try:
                if language == "si":
                    question = f"'{concept}' ගැන කෙටියෙන් පැහැදිලි කරන්න."
                elif language == "ta":
                    question = f"'{concept}' பற்றி சுருக்கமாக விளக்கவும்."
                else:
                    question = f"Briefly explain '{concept}'."
                
                # Find relevant sentence containing the concept
                relevant_sentence = next((s for s in sentences if concept.lower() in s.lower()), sentences[0])
                
                exercises.append({
                    "type": "short_answer",
                    "instruction": "Answer the following question:",
                    "exercise_text": question,
                    "answer": relevant_sentence,
                    "difficulty": difficulty
                })
            except Exception as e:
                continue
        
        # Exercise Type 4: Matching (if enough concepts)
        if len(key_concepts) >= 4:
            try:
                # Create matching pairs
                concepts_subset = key_concepts[:4]
                definitions = []
                
                for concept in concepts_subset:
                    # Find sentence containing the concept
                    definition = next((s for s in sentences if concept.lower() in s.lower()), f"Related to {concept}")
                    definitions.append(definition[:100] + "..." if len(definition) > 100 else definition)
                
                if language == "si":
                    instruction = "පහත සංකල්ප සහ අර්ථ දැක්වීම් ගලපන්න:"
                elif language == "ta":
                    instruction = "பின்வரும் கருத்துகள் மற்றும் வரையறைகளை பொருத்தவும்:"
                else:
                    instruction = "Match the following concepts with their definitions:"
                
                exercises.append({
                    "type": "matching",
                    "instruction": instruction,
                    "concepts": concepts_subset,
                    "definitions": definitions,
                    "answer": dict(zip(concepts_subset, definitions)),
                    "difficulty": difficulty
                })
            except Exception as e:
                logger.warning(f"Error creating matching exercise: {str(e)}")
        
        logger.info(f"Generated {len(exercises)} exercises")
        return exercises
        
    except Exception as e:
        logger.error(f"Error generating exercises: {str(e)}")
        return []

def get_card_types(card_types: Optional[List[str]] = Form(None)):
    return card_types or ["flashcard"]

@app.post("/process-file")
async def process_file(
    file: UploadFile = File(...),
    language: str = Form("en"),
    card_types: List[str] = Depends(get_card_types),
    difficulty: str = Form("beginner")
):
    """
    Process uploaded document and generate flashcards, exercises, and quizzes.
    
    Args:
        file: The uploaded file (PDF, DOCX, or PPTX)
        language: The language of the document (en, si, ta)
        card_types: List of card types to generate (flashcard, exercise, quiz)
        difficulty: The difficulty level (beginner, intermediate, advanced)
    
    Returns:
        JSON response containing generated content
    """
    logger.info(f"Processing file: {file.filename}")
    logger.debug(f"Parameters: language={language}, card_types={card_types}, difficulty={difficulty}")
    
    try:
        # Validate parameters
        if language not in ["en", "si", "ta"]:
            raise HTTPException(status_code=400, detail="Unsupported language")
        
        card_types_list = [ct for ct in card_types if ct in ["flashcard", "exercise", "quiz"]]
        if not card_types_list:
            card_types_list = ["flashcard"]
        
        if difficulty not in ["beginner", "intermediate", "advanced"]:
            difficulty = "beginner"
        
        # Read and process file
        content = await file.read()
        file_extension = file.filename.split(".")[-1].lower()
        
        # Extract text based on file type
        if file_extension == "pdf":
            text = extract_text_from_pdf(content)
        elif file_extension == "docx":
            text = extract_text_from_docx(content)
        elif file_extension == "pptx":
            text = extract_text_from_pptx(content)
        else:
            raise HTTPException(status_code=400, detail="Unsupported file type")
        
        if not text.strip():
            raise HTTPException(status_code=400, detail="No text content found in the document")
        
        # Generate requested content types
        generated_content = {}
        
        if "flashcard" in card_types_list:
            logger.info("Generating flashcards...")
            flashcards = generate_flashcards(text, language, difficulty)
            generated_content["flashcards"] = flashcards
        
        if "quiz" in card_types_list:
            logger.info("Generating quizzes...")
            quizzes = generate_quizzes(text, language, difficulty)
            generated_content["quizzes"] = quizzes
        
        if "exercise" in card_types_list:
            logger.info("Generating exercises...")
            exercises = generate_exercises(text, language, difficulty)
            generated_content["exercises"] = exercises
        
        logger.info(f"Successfully processed {file.filename}. Generated: {list(generated_content.keys())}")
        return {"generated_content": generated_content}
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Unexpected error processing {file.filename}: {str(e)}", exc_info=True)
        raise HTTPException(status_code=500, detail="Internal server error")

# Add health check endpoint
@app.get("/health")
async def health_check():
    """Health check endpoint."""
    return {
        "status": "healthy",
        "model": current_model_name,
        "timestamp": datetime.now().isoformat()
    }

# Add middleware to log requests
@app.middleware("http")
async def log_requests(request, call_next):
    start_time = datetime.now()
    logger.info(f"Request: {request.method} {request.url}")
    
    try:
        response = await call_next(request)
        process_time = (datetime.now() - start_time).total_seconds()
        logger.info(f"Response: {response.status_code} - {process_time:.3f}s")
        return response
    except Exception as e:
        logger.error(f"Request failed: {str(e)}", exc_info=True)
        raise

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8001)